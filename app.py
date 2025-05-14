
import streamlit as st
import whisper
import openai
from pptx import Presentation
from tempfile import NamedTemporaryFile
import subprocess
import os
import io

# Azure setup (use secrets in Streamlit Cloud or local .streamlit/secrets.toml)
AZURE_API_KEY = st.secrets["AZURE_API_KEY"]
AZURE_ENDPOINT = st.secrets["AZURE_ENDPOINT"]
AZURE_DEPLOYMENT = st.secrets["AZURE_DEPLOYMENT"]

openai.api_key = AZURE_API_KEY
openai.api_base = AZURE_ENDPOINT
openai.api_type = "azure"
openai.api_version = "2025-03-01-preview"

# Convert MP4 to WAV
def convert_mp4_to_wav(mp4_bytes, output_path):
    with NamedTemporaryFile(delete=False, suffix=".mp4") as temp_file:
        temp_file.write(mp4_bytes)
        temp_file.flush()
        subprocess.call(['ffmpeg', '-i', temp_file.name, output_path])
    os.remove(temp_file.name)

# Transcribe Audio
def transcribe_audio(wav_path):
    model = whisper.load_model("base")
    result = model.transcribe(wav_path)
    return result['text']

# Extract Text from PPTX
def extract_text_from_pptx(pptx_file):
    prs = Presentation(pptx_file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

# Call LLM for QC
def quality_check(video_text, ppt_text):
    prompt = f"""
    Perform a quality check between the video transcript and the PPT content.

    - Check if key concepts are covered
    - Check alignment between spoken and visual content
    - Check clarity of learning outcomes
    - Check for code walkthroughs

    Return a checklist with ✅ or ❌ and reasoning.

    Transcript:
    {video_text}

    Slides:
    {ppt_text}
    """

    response = openai.ChatCompletion.create(
        engine=AZURE_DEPLOYMENT,
        messages=[{"role": "user", "content": prompt}]
    )
    return response.choices[0].message["content"]

# Streamlit UI
st.title("🎧 Recording Quality Check with LLM")

mp4_file = st.file_uploader("Upload MP4 Video", type=["mp4"])
pptx_file = st.file_uploader("Upload PPTX Slides", type=["pptx"])

if st.button("Run Quality Check") and mp4_file and pptx_file:
    with st.spinner("Processing..."):
        convert_mp4_to_wav(mp4_file.read(), "audio.wav")
        transcript = transcribe_audio("audio.wav")
        ppt_text = extract_text_from_pptx(pptx_file)
        report = quality_check(transcript, ppt_text)

        st.success("✅ Quality Check Complete")
        st.text_area("QC Report", report, height=400)

        # Enable download
        report_file = io.StringIO(report)
        st.download_button(
            label="📥 Download QC Report",
            data=report_file,
            file_name="QC_Report.txt",
            mime="text/plain"
        )
